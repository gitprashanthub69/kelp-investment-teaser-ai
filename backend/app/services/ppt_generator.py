"""
PPT Generator - Sample Output Style
Dense professional layout with logos, certifications, assets grid
4 Slides: Title + Business Overview + Financials + Investment Highlights
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from typing import Dict, Any, List, Optional
import os
import random
import requests
import tempfile
import io
import urllib.parse
try:
    from PIL import Image
except ImportError:
    Image = None

# For advanced logo discovery
try:
    from app.services.scraper import ScraperService
except ImportError:
    ScraperService = None


class PPTGenerator:
    """Sample Output Style PPT Generator"""
    
    # Professional Colors
    NAVY = RGBColor(24, 40, 72)
    TEAL = RGBColor(0, 128, 128)
    LIGHT_TEAL = RGBColor(0, 176, 204)
    ORANGE = RGBColor(255, 140, 0)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    MED_GRAY = RGBColor(220, 220, 220)
    DARK_GRAY = RGBColor(80, 80, 80)
    TEXT_DARK = RGBColor(33, 33, 33)
    GREEN = RGBColor(40, 167, 69)
    
    CODENAMES = ["Apex", "Stellar", "Horizon", "Summit", "Pinnacle", "Atlas", "Titan", "Phoenix", "Orion", "Eclipse", "Nova", "Zenith"]
    
    # Common certifications
    CERTIFICATIONS = {
        "pharma": ["WHO-GMP", "FDA", "EU-GMP", "ISO 9001", "ISO 14001"],
        "manufacturing": ["ISO 9001", "IATF 16949", "ISO 14001", "OHSAS 18001"],
        "food": ["FSSAI", "ISO 22000", "HACCP", "BRC", "FSSC 22000"],
        "general": ["ISO 9001", "ISO 14001", "SEDEX", "SA 8000"]
    }

    def __init__(self, output_path: str, data: Dict[str, Any]):
        self.output_path = output_path
        self.data = data
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.codename = f"Project {random.choice(self.CODENAMES)}"
        self.sector = data.get("sector", "General Business")
        self.temp_images = []  # Track temp image files for cleanup
        self.logo_cache = {}  # Cache downloaded logos

    def generate(self) -> str:
        try:
            self._slide_1_title()
            self._slide_2_business_overview()
            self._slide_3_financial_performance()
            self._slide_4_investment_highlights()
            self.prs.save(self.output_path)
            print(f"[PPTGenerator] Saved: {self.output_path}")
            return self.output_path
        except Exception as e:
            print(f"[PPTGenerator] Error: {e}")
            raise

    # ==================== SLIDE 1: TITLE ====================
    def _slide_1_title(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # Navy background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.NAVY
        
        # Logo placeholder
        self._add_box(slide, 0.4, 0.4, 2.2, 0.7, self.WHITE)
        self._add_text(slide, 0.5, 0.48, 2.0, 0.55, "KELP", 26, self.NAVY, bold=True)
        
        # Project name (significantly larger)
        self._add_text(slide, 0.4, 2.5, 9.0, 1.6, self.codename, 72, self.WHITE, bold=True)
        
        # Sector badge (larger)
        self._add_box(slide, 0.4, 4.3, 5.0, 0.7, self.LIGHT_TEAL)
        self._add_text(slide, 0.5, 4.38, 4.8, 0.55, self.sector.upper(), 22, self.WHITE, bold=True)
        
        # Tagline (larger)
        self._add_text(slide, 0.4, 5.3, 9.0, 0.8, self.data.get("tagline", "UNLEASHING STRATEGIC VALUE"), 16, self.WHITE, italic=True)
        self._add_text(slide, 0.4, 5.7, 7.0, 0.6, "Investment Teaser | Strictly Confidential", 16, self.WHITE)
        
        # Right side visual
        self._add_box(slide, 9.0, 0.8, 4.0, 6.0, RGBColor(40, 60, 100))
        self._add_text(slide, 9.2, 3.0, 3.6, 1.0, "Company\nVisual", 28, RGBColor(80, 100, 140), center=True)
        
        # Footer
        self._add_text(slide, 0.4, 7.0, 12.5, 0.35, "kelpglobal.com", 12, self.WHITE, center=True)

    # ==================== SLIDE 2: BUSINESS OVERVIEW ====================
    def _slide_2_business_overview(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.LIGHT_GRAY
        
        # Header
        self._add_header(slide, "Business Overview")
        
        narrative = self.data.get("narrative", {}).get("slide_1", {})
        
        # === LEFT COLUMN (0.3 - 6.4) ===
        
        # Business Description (compact but larger font)
        self._add_section_header(slide, 0.3, 1.1, "COMPANY SNAPSHOT")
        biz_desc = narrative.get("biz_desc", self._default_biz_desc())[:350]
        self._add_content_box(slide, 0.3, 1.35, 6.1, 1.25, biz_desc, 11)
        
        # Key Select Customers
        self._add_section_header(slide, 0.3, 2.8, "KEY SELECT CUSTOMERS")
        customers = narrative.get("customers", ["Customer A", "Customer B", "Customer C", "Customer D", "Customer E", "Customer F"])
        self._add_customer_grid(slide, 0.3, 3.05, 6.1, customers[:6])
        
        # Company at a Glance (Assets Grid)
        self._add_section_header(slide, 0.3, 4.35, "COMPANY AT A GLANCE")
        assets = narrative.get("assets", self._default_assets())
        self._add_assets_grid(slide, 0.3, 4.6, 6.1, assets)
        
        # Certifications
        self._add_section_header(slide, 0.3, 5.85, "CERTIFICATIONS & ACCREDITATIONS")
        certs = narrative.get("certifications", self._get_sector_certs())
        self._add_certification_row(slide, 0.3, 6.1, 6.1, certs)
        
        # === RIGHT COLUMN (6.6 - 12.9) ===
        
        # Product Portfolio
        self._add_section_header(slide, 6.6, 1.1, "PRODUCT PORTFOLIO")
        portfolio = narrative.get("product_portfolio", self._default_portfolio())
        self._add_product_list(slide, 6.6, 1.35, 6.3, portfolio[:4])
        
        # Key Applications
        self._add_section_header(slide, 6.6, 3.9, "KEY APPLICATIONS / END MARKETS")
        apps = narrative.get("applications", self._default_applications())
        self._add_applications_grid(slide, 6.6, 4.15, 6.3, apps[:4]) # Limit to 4 for clarity
        
        # Revenue Split
        self._add_section_header(slide, 6.6, 5.7, "REVENUE MIX")
        rev_split = narrative.get("revenue_split", {"Domestic": 65, "Export": 35})
        self._add_revenue_bars(slide, 6.6, 5.95, 6.3, rev_split)
        
        self._add_footer(slide)

    # ==================== SLIDE 3: FINANCIAL PERFORMANCE ====================
    def _slide_3_financial_performance(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.LIGHT_GRAY
        
        self._add_header(slide, "Financial Performance & Global Presence")
        
        financials = self.data.get("financials", {})
        narrative = self.data.get("narrative", {}).get("slide_3", {})
        kpis = self.data.get("kpis", {})
        
        # === TOP ROW: Key Metrics ===
        self._add_section_header(slide, 0.3, 1.1, "KEY FINANCIAL METRICS")
        metrics = self._extract_financial_metrics(financials, kpis)
        self._add_metrics_row(slide, 0.3, 1.35, 8.5, metrics[:5])
        
        # === LEFT: Revenue Chart ===
        self._add_section_header(slide, 0.3, 2.4, "REVENUE & EBITDA TREND")
        years = financials.get("years", ["FY21", "FY22", "FY23", "FY24"])
        revenue = financials.get("revenue", [100, 125, 160, 200])
        if len(years) >= 2:
            self._add_combo_chart(slide, 0.3, 2.65, 5.5, 2.2, years, revenue, financials.get("ebitda", []))
        
        # === RIGHT: Global Presence ===
        self._add_section_header(slide, 6.0, 2.4, "GLOBAL PRESENCE")
        global_text = narrative.get("global_reach", "Pan-India presence with exports to 20+ countries across Asia, Europe, and Americas.")
        self._add_content_box(slide, 6.0, 2.65, 6.8, 1.0, global_text, 9)
        
        # Export Markets
        exports = narrative.get("export_markets", ["USA", "Europe", "Middle East", "SEA", "LATAM"])
        self._add_market_badges(slide, 6.0, 3.75, 6.8, exports[:5])
        
        # === BOTTOM LEFT: Upcoming Facilities ===
        self._add_section_header(slide, 0.3, 5.0, "UPCOMING CAPACITY / FACILITIES")
        upcoming = narrative.get("upcoming_facilities", [
            "New manufacturing unit - 50,000 sq ft (FY25)",
            "R&D center expansion (Q2 FY25)",
            "Warehouse automation project"
        ])
        # Wrap in box for Phase 4
        self._add_content_box(slide, 0.3, 5.25, 5.5, 1.25, "\n• " + "\n• ".join(upcoming[:3]), 11)
        
        # === BOTTOM RIGHT: Assumptions ===
        self._add_section_header(slide, 6.0, 5.0, "KEY ASSUMPTIONS")
        assumptions = narrative.get("assumptions", [
            "Financials based on audited statements",
            "Projections assume current market conditions",
            "Export figures in USD converted at prevailing rates"
        ])
        # Wrap in box for Phase 4
        self._add_content_box(slide, 6.0, 5.25, 6.8, 1.25, "\n• " + "\n• ".join(assumptions[:3]), 11)
        
        self._add_footer(slide)

    # ==================== SLIDE 4: INVESTMENT HIGHLIGHTS ====================
    def _slide_4_investment_highlights(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.LIGHT_GRAY
        
        self._add_header(slide, "Investment Highlights")
        
        narrative = self.data.get("narrative", {}).get("slide_3", {})
        
        # Investment Highlights (6 cards in 2x3 grid)
        highlights = narrative.get("investment_highlights", self._default_investment_highlights())
        
        card_w = 4.0
        card_h = 1.85  # Increased to prevent content overflow with larger fonts
        start_x = 0.4
        start_y = 1.2
        gap_x = 0.25
        gap_y = 0.15
        
        for i, hl in enumerate(highlights[:6]):
            row = i // 3
            col = i % 3
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)
            self._add_investment_card(slide, x, y, card_w, card_h, i+1, hl.get("title", ""), hl.get("desc", ""))
        
        # Bottom Summary Box
        self._add_box(slide, 0.4, 5.2, 12.5, 1.6, self.NAVY)
        self._add_text(slide, 0.6, 5.3, 12.1, 0.45, "WHY INVEST?", 22, self.LIGHT_TEAL, bold=True)
        
        why_invest = narrative.get("why_invest", 
            "Strong market position with consistent growth trajectory, experienced management team, diversified revenue streams, and clear expansion roadmap make this an attractive investment opportunity with significant upside potential.")
        self._add_text(slide, 0.6, 5.75, 12.1, 0.95, why_invest, 15, self.WHITE)
        
        self._add_footer(slide)

    # ==================== HELPER METHODS ====================
    
    def _add_header(self, slide, title: str):
        # Header bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.NAVY
        bar.line.fill.background()
        
        # Logo
        self._add_text(slide, 0.3, 0.25, 2.0, 0.4, "KELP", 26, self.WHITE, bold=True)
        
        # Title
        self._add_text(slide, 3.0, 0.2, 7.5, 0.6, title, 32, self.WHITE, bold=True, center=True)
        
        # Codename
        self._add_text(slide, 10.5, 0.25, 2.5, 0.4, self.codename, 16, self.LIGHT_TEAL)

    def _add_footer(self, slide):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.33), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.MED_GRAY
        line.line.fill.background()
        
        self._add_text(slide, 0.3, 7.15, 12.7, 0.25, "Strictly Private & Confidential | Prepared by Kelp M&A Team", 9, self.DARK_GRAY, center=True)

    def _add_section_header(self, slide, x, y, text):
        self._add_text(slide, x, y, 6.0, 0.3, text, 14, self.TEAL, bold=True)

    def _add_box(self, slide, x, y, w, h, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_content_box(self, slide, x, y, w, h, text, font_size):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = self.WHITE
        box.line.color.rgb = self.MED_GRAY
        box.line.width = Pt(0.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.color.rgb = self.TEXT_DARK
        p.line_spacing = 1.15
        p.alignment = PP_ALIGN.LEFT
        return box

    def _add_text(self, slide, x, y, w, h, text, size, color, bold=False, center=False, italic=False, v_center=True):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        
        # Better orientation: vertical centering and consistent margins
        if v_center:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(5)
        tf.margin_right = Pt(5)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.line_spacing = 1.1
        if center:
            p.alignment = PP_ALIGN.CENTER
        return box

    def _add_customer_grid(self, slide, x, y, w, customers):
        """Add customer logos/names in a grid - attempts to fetch real logos"""
        box_w = (w - 0.25) / 3
        box_h = 0.5
        
        for i, cust in enumerate(customers[:6]):
            row = i // 3
            col = i % 3
            bx = x + col * (box_w + 0.08)
            by = y + row * (box_h + 0.1)
            
            # Create box background
            box = self._add_box(slide, bx, by, box_w, box_h, self.WHITE)
            box.line.color.rgb = self.MED_GRAY
            box.line.width = Pt(0.5)
            
            # Try to add logo image
            logo_added = self._try_add_logo(slide, str(cust), bx + 0.05, by + 0.08, box_w - 0.1, box_h - 0.16)
            
            if not logo_added:
                # If no logo, show text name (centered and oriented, larger font for Phase 4)
                self._add_text(slide, bx, by, box_w, box_h, str(cust)[:25], 14, self.NAVY, bold=True, center=True, v_center=True)

    def _add_assets_grid(self, slide, x, y, w, assets):
        """Add company assets (facilities, R&D, etc)"""
        box_w = (w - 0.4) / 4
        box_h = 1.0
        
        for i, asset in enumerate(assets[:4]):
            bx = x + i * (box_w + 0.1)
            
            box = self._add_box(slide, bx, y, box_w, box_h, self.NAVY)
            
            # Value (larger font)
            self._add_text(slide, bx, y + 0.1, box_w, 0.45, str(asset.get("value", "-")), 26, self.WHITE, bold=True, center=True)
            # Label (larger font)
            self._add_text(slide, bx, y + 0.6, box_w, 0.4, str(asset.get("label", "")), 11, self.LIGHT_TEAL, center=True)

    def _add_certification_row(self, slide, x, y, w, certs):
        """Add certification badges with icons"""
        badge_w = (w - 0.4) / 5
        
        for i, cert in enumerate(certs[:5]):
            bx = x + i * (badge_w + 0.08)
            
            badge = self._add_box(slide, bx, y, badge_w, 0.5, self.WHITE)
            badge.line.color.rgb = self.TEAL
            badge.line.width = Pt(1.5)
            
            # Try to add cert icon/logo
            cert_logo_added = self._try_add_cert_icon(slide, str(cert), bx + 0.08, y + 0.08, badge_w - 0.16, 0.35)
            
            if not cert_logo_added:
                # Cert name (centered and oriented)
                self._add_text(slide, bx, y, badge_w, 0.5, str(cert), 11, self.TEAL, bold=True, center=True, v_center=True)

    def _add_product_list(self, slide, x, y, w, products):
        """Add product portfolio items"""
        item_h = 0.6
        
        for i, prod in enumerate(products[:4]):
            py = y + i * (item_h + 0.05)
            
            # Colored left bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(py), Inches(0.1), Inches(item_h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.TEAL
            bar.line.fill.background()
            
            # Content box
            box = self._add_box(slide, x + 0.1, py, w - 0.1, item_h, self.WHITE)
            box.line.color.rgb = self.MED_GRAY
            box.line.width = Pt(0.5)
            
            # Title (larger font)
            self._add_text(slide, x + 0.2, py + 0.08, w - 0.35, 0.3, str(prod.get("category", ""))[:28], 14, self.NAVY, bold=True)
            # Details (larger font)
            self._add_text(slide, x + 0.2, py + 0.35, w - 0.35, 0.3, str(prod.get("details", ""))[:65], 12, self.DARK_GRAY)

    def _add_applications_grid(self, slide, x, y, w, apps):
        """Add applications with percentage bars"""
        item_h = 0.35
        
        for i, app in enumerate(apps[:4]):
            ay = y + i * (item_h + 0.08)
            
            # Background
            bg = self._add_box(slide, x, ay, w, item_h, self.WHITE)
            bg.line.color.rgb = self.MED_GRAY
            bg.line.width = Pt(0.5)
            
            # Industry name
            self._add_text(slide, x + 0.1, ay + 0.05, w * 0.65, 0.25, str(app.get("industry", "")), 11, self.TEXT_DARK)
            
            # Share badge
            share = str(app.get("share", ""))
            badge = self._add_box(slide, x + w - 0.75, ay + 0.03, 0.7, 0.28, self.NAVY)
            self._add_text(slide, x + w - 0.75, ay + 0.03, 0.7, 0.28, share, 10, self.WHITE, bold=True, center=True, v_center=True)

    def _add_revenue_bars(self, slide, x, y, w, rev_split):
        """Add horizontal revenue split bars"""
        bar_h = 0.3
        total = sum(rev_split.values()) or 1
        
        current_x = x
        colors = [self.NAVY, self.TEAL, self.LIGHT_TEAL, self.ORANGE]
        
        for i, (label, value) in enumerate(rev_split.items()):
            bar_w = (value / total) * w
            
            bar = self._add_box(slide, current_x, y, bar_w - 0.02, bar_h, colors[i % len(colors)])
            self._add_text(slide, current_x, y + 0.05, bar_w - 0.02, 0.2, f"{label}: {value}%", 8, self.WHITE, center=True)
            
            current_x += bar_w

    def _add_metrics_row(self, slide, x, y, w, metrics):
        """Add horizontal metrics row"""
        box_w = (w - 0.4) / len(metrics) if metrics else 1.5
        
        for i, m in enumerate(metrics):
            mx = x + i * (box_w + 0.08)
            
            box = self._add_box(slide, mx, y, box_w, 0.75, self.WHITE)
            box.line.color.rgb = self.NAVY
            box.line.width = Pt(1)
            
            # Value
            self._add_text(slide, mx, y + 0.1, box_w, 0.35, str(m.get("value", "-")), 16, self.NAVY, bold=True, center=True)
            # Label
            self._add_text(slide, mx, y + 0.45, box_w, 0.25, str(m.get("label", "")), 7, self.DARK_GRAY, center=True)

    def _add_combo_chart(self, slide, x, y, w, h, years, revenue, ebitda):
        """Add bar chart for revenue"""
        try:
            chart_data = CategoryChartData()
            chart_data.categories = years
            chart_data.add_series('Revenue', revenue)
            if ebitda:
                chart_data.add_series('EBITDA', ebitda)
            
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = Pt(7)
        except Exception as e:
            print(f"[PPTGenerator] Chart error: {e}")

    def _add_market_badges(self, slide, x, y, w, markets):
        """Add export market badges"""
        badge_w = (w - 0.4) / len(markets) if markets else 1.2
        
        for i, market in enumerate(markets):
            mx = x + i * (badge_w + 0.08)
            badge = self._add_box(slide, mx, y, badge_w, 0.4, self.NAVY)
            self._add_text(slide, mx, y + 0.08, badge_w, 0.28, str(market), 10, self.WHITE, center=True)

    def _add_bullet_list(self, slide, x, y, w, items, font_size, italic=False):
        """Add bullet list"""
        for i, item in enumerate(items):
            iy = y + i * 0.35
            self._add_text(slide, x, iy, 0.15, 0.25, "•", font_size, self.TEAL, bold=True)
            self._add_text(slide, x + 0.18, iy, w - 0.2, 0.3, str(item), font_size, self.TEXT_DARK, italic=italic)

    def _add_investment_card(self, slide, x, y, w, h, num, title, desc):
        """Add investment highlight card"""
        # Card background
        card = self._add_box(slide, x, y, w, h, self.WHITE)
        card.line.color.rgb = self.MED_GRAY
        card.line.width = Pt(0.5)
        
        # Number circle (larger)
        circle_size = 0.5
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.12), Inches(y + 0.12), Inches(circle_size), Inches(circle_size))
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.NAVY
        circle.line.fill.background()
        self._add_text(slide, x + 0.12, y + 0.2, circle_size, 0.32, str(num), 16, self.WHITE, bold=True, center=True)
        
        # Title (larger font)
        self._add_text(slide, x + 0.7, y + 0.2, w - 0.85, 0.4, title, 16, self.NAVY, bold=True)
        
        # Description (larger font)
        self._add_text(slide, x + 0.15, y + 0.75, w - 0.3, h - 0.85, desc, 14, self.DARK_GRAY, v_center=False)

    # ==================== DATA EXTRACTION ====================
    
    def _extract_financial_metrics(self, financials, kpis):
        """Extract key financial metrics"""
        metrics = []
        
        revenue = financials.get("revenue", [])
        ebitda = financials.get("ebitda", [])
        
        if revenue:
            metrics.append({"label": "Revenue (FY24)", "value": f"₹{revenue[-1]} Cr"})
            cagr = self._calc_cagr(revenue)
            if cagr:
                metrics.append({"label": "Revenue CAGR", "value": f"{cagr}%"})
        
        if ebitda and revenue:
            margin = round((ebitda[-1] / revenue[-1]) * 100) if revenue[-1] > 0 else 0
            metrics.append({"label": "EBITDA Margin", "value": f"{margin}%"})
        
        metrics.append({"label": "RoCE", "value": kpis.get("roce", "18%")})
        metrics.append({"label": "Debt", "value": kpis.get("debt", "Low")})
        
        return metrics

    def _calc_cagr(self, series):
        try:
            if len(series) < 2:
                return None
            valid = [x for x in series if x and float(x) > 0]
            if len(valid) < 2:
                return None
            return int(round(((float(valid[-1]) / float(valid[0])) ** (1 / (len(valid) - 1)) - 1) * 100))
        except:
            return None

    def _get_sector_certs(self):
        sector_lower = self.sector.lower()
        if "pharma" in sector_lower or "health" in sector_lower:
            return self.CERTIFICATIONS["pharma"]
        elif "food" in sector_lower or "fmcg" in sector_lower or "consumer" in sector_lower:
            return self.CERTIFICATIONS["food"]
        elif "manufactur" in sector_lower:
            return self.CERTIFICATIONS["manufacturing"]
        return self.CERTIFICATIONS["general"]

    # ==================== DEFAULTS ====================
    
    def _default_biz_desc(self):
        return f"The Company is a leading player in the {self.sector} sector with diversified operations, strong customer relationships, and established market position. Built on operational excellence with an experienced management team driving consistent growth."

    def _default_assets(self):
        return [
            {"label": "Manufacturing\nUnits", "value": "3"},
            {"label": "R&D\nCenters", "value": "1"},
            {"label": "Employees", "value": "500+"},
            {"label": "Years in\nBusiness", "value": "15+"}
        ]

    def _default_portfolio(self):
        return [
            {"category": "Core Products", "details": "Primary revenue driver with market leadership"},
            {"category": "Specialty Range", "details": "High-margin differentiated products"},
            {"category": "Contract Services", "details": "B2B manufacturing and export services"},
            {"category": "New Launches", "details": "Pipeline products for growth markets"}
        ]

    def _default_applications(self):
        return [
            {"industry": "Primary Sector", "share": "40%"},
            {"industry": "Secondary Sector", "share": "25%"},
            {"industry": "Tertiary Sector", "share": "20%"},
            {"industry": "Others", "share": "15%"}
        ]

    def _default_investment_highlights(self):
        return [
            {"title": "Market Leadership", "desc": "Established brand with strong market position and customer loyalty in core segments"},
            {"title": "Diversified Portfolio", "desc": "Balanced product mix reducing concentration risk with multiple revenue streams"},
            {"title": "Scalable Operations", "desc": "Modern infrastructure with expansion capacity to support 2x growth"},
            {"title": "Experienced Team", "desc": "Strong management with 20+ years average industry experience"},
            {"title": "Growth Trajectory", "desc": "Consistent 20%+ CAGR with improving margins and operational efficiency"},
            {"title": "Clear Roadmap", "desc": "Well-defined expansion plans including new facilities and market entry"}
        ]

    # ==================== LOGO/IMAGE METHODS ====================
    
    def _try_add_logo(self, slide, company_name: str, x: float, y: float, w: float, h: float) -> bool:
        """Try to add a company logo using Clearbit Logo API with multiple domain variations"""
        try:
            # Clean company name
            raw_name = company_name.lower().strip()
            clean_name = raw_name
            # Remove common business suffixes to get the core name
            suffixes = [
                " ltd", " limited", " pvt", " private", " inc", " corp", " llc", " llp", " co",
                " solutions", " technologies", " systems", " services", " group", " holdings"
            ]
            for suffix in suffixes:
                clean_name = clean_name.replace(suffix, "")
            clean_name = clean_name.strip().replace(" ", "")
            
            if not clean_name or len(clean_name) < 2:
                return self._add_letter_logo(slide, raw_name, x, y, w, h)
            
            # Check cache first
            if clean_name in self.logo_cache:
                img_path = self.logo_cache[clean_name]
                if img_path and os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
                    return True
                if img_path is None: # Known failure
                     return self._add_letter_logo(slide, raw_name, x, y, w, h)
            
            # Try multiple TLDs
            tlds = [".com", ".in", ".co.in", ".net", ".org"]
            for tld in tlds:
                logo_url = f"https://logo.clearbit.com/{clean_name}{tld}"
                try:
                    response = requests.get(logo_url, timeout=3)
                    if response.status_code == 200 and 'image' in response.headers.get('content-type', ''):
                        # Save to temp file
                        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                        temp_file.write(response.content)
                        temp_file.close()
                        self.temp_images.append(temp_file.name)
                        self.logo_cache[clean_name] = temp_file.name
                        
                        # Add to slide
                        slide.shapes.add_picture(temp_file.name, Inches(x), Inches(y), Inches(w), Inches(h))
                        print(f"[PPTGenerator] Added logo for: {company_name} using {tld}")
                        return True
                except:
                    continue
            
            # --- PHASE 4: WIKIPEDIA & GOOGLE SEARCH logo discovery ---
            if ScraperService and google_search:
                try:
                    # Search specifically for Wikipedia or high-res logos
                    q = f"{company_name} logo site:wikipedia.org OR site:wikimedia.org"
                    results = list(google_search(q, num_results=3, lang="en"))
                    for url in results:
                        # Try to find an image file on the page or use the page to find the official logo
                        if "wikipedia.org" in url or "wikimedia.org" in url:
                            resp = requests.get(url, timeout=5)
                            if resp.status_code == 200:
                                soup = BeautifulSoup(resp.text, 'html.parser')
                                # Look for the main logo in the infobox
                                infobox = soup.find('table', class_='infobox')
                                if infobox:
                                    img = infobox.find('img')
                                    if img and img.get('src'):
                                        img_url = "https:" + img['src'] if img['src'].startswith("//") else img['src']
                                        # download and add
                                        img_resp = requests.get(img_url, timeout=5)
                                        if img_resp.status_code == 200:
                                            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                                            temp_file.write(img_resp.content)
                                            temp_file.close()
                                            self.temp_images.append(temp_file.name)
                                            slide.shapes.add_picture(temp_file.name, Inches(x), Inches(y), Inches(w), Inches(h))
                                            print(f"[PPTGenerator] Added logo for: {company_name} from Wikipedia")
                                            return True
                except Exception as e:
                    print(f"[PPTGenerator] Wikipedia discovery failed for {company_name}: {e}")

            # Cache failed attempt
            self.logo_cache[clean_name] = None
            return self._add_letter_logo(slide, raw_name, x, y, w, h)
            
        except Exception as e:
            print(f"[PPTGenerator] Logo fetch failed for {company_name}: {e}")
            return self._add_letter_logo(slide, raw_name, x, y, w, h)

    def _add_letter_logo(self, slide, name: str, x: float, y: float, w: float, h: float) -> bool:
        """Create a professional styled 'Letter Logo' fallback"""
        try:
            letter = name.strip()[0].upper() if name.strip() else "?"
            # Background badge
            badge_size = min(w, h) * 0.9
            bx = x + (w - badge_size) / 2
            by = y + (h - badge_size) / 2
            
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx), Inches(by), Inches(badge_size), Inches(badge_size))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.NAVY
            badge.line.color.rgb = self.LIGHT_TEAL
            badge.line.width = Pt(1)
            
            # Letter text
            self._add_text(slide, bx, by + (badge_size * 0.15), badge_size, badge_size * 0.7, letter, 18, self.WHITE, bold=True, center=True)
            return True
        except:
            return False
    
    def _try_add_cert_icon(self, slide, cert_name: str, x: float, y: float, w: float, h: float) -> bool:
        """Try to add a certification icon using a predefined domain mapping"""
        cert_domains = {
            "ISO 9001": "iso.org",
            "ISO 14001": "iso.org",
            "ISO 22000": "iso.org",
            "ISO 27001": "iso.org",
            "WHO-GMP": "who.int",
            "FDA": "fda.gov",
            "FSSAI": "fssai.gov.in",
            "HACCP": "haccp.com",
            "BRC": "brcgs.com",
            "IATF 16949": "iatfglobaloversight.org",
            "OHSAS 18001": "iso.org",
            "REACH": "echa.europa.eu",
            "GDPR": "gdpr.eu",
            "SOC 2": "aicpa.org",
            "SEDEX": "sedex.com",
            "SA 8000": "sa-intl.org"
        }
        
        # Direct match or partial match
        name = cert_name.upper().strip()
        domain = None
        for key, val in cert_domains.items():
            if key in name or name in key:
                domain = val
                break
        
        if not domain:
            return False
            
        try:
            logo_url = f"https://logo.clearbit.com/{domain}"
            response = requests.get(logo_url, timeout=3)
            if response.status_code == 200 and 'image' in response.headers.get('content-type', ''):
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                temp_file.write(response.content)
                temp_file.close()
                self.temp_images.append(temp_file.name)
                slide.shapes.add_picture(temp_file.name, Inches(x), Inches(y), Inches(w), Inches(h))
                return True
        except:
            pass
            
        return False
    
    def _cleanup_temp_images(self):
        """Clean up temporary image files"""
        for temp_path in self.temp_images:
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
            except:
                pass
        self.temp_images = []
